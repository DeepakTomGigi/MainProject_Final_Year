from pptx import Presentation
from pptx.util import Inches, Pt

def fits_on_slide(text, max_chars=500):
    return len(text) <= max_chars

# Function to process the summary and generate the presentation
def generate_presentation(summary_text):
    # Ensure input is not empty
    if not summary_text.strip():
        print("Error: No summary text provided.")
        return

    # Parse the summary text into sections
    sections = {}
    current_section = None
    for line in summary_text.splitlines():
        line = line.strip()
        if line.startswith("[") and line.endswith("]"):
            current_section = line[1:-1].lower()
            sections[current_section] = []
        elif current_section and line:
            sections[current_section].append(line)

    # Validate sections
    if "title" not in sections or not sections["title"]:
        print("Error: [Title] section is missing or empty.")
        return
    if "content" not in sections or not sections["content"]:
        print("Error: [Content] section is missing or empty.")
        return

    # Create a PowerPoint presentation
    presentation = Presentation()

    # Add title slide
    title_slide_layout = presentation.slide_layouts[0]
    slide = presentation.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = sections["title"][0]
    subtitle = sections["title"][1] if len(sections["title"]) > 1 else "Generated Presentation"
    slide.placeholders[1].text = subtitle

    # Add content slides
    content_slide_layout = presentation.slide_layouts[1]
    current_slide = None
    current_text = ""

    for line in sections["content"]:
        if not line.startswith("-"):  # New heading -> start a new slide
            if current_slide:  # Save any previous slide content
                text_frame = current_slide.shapes[1].text_frame
                text_frame.text = current_text.strip()
            current_slide = presentation.slides.add_slide(content_slide_layout)
            current_slide.shapes.title.text = line
            current_text = ""
        else:  # Bullet point
            bullet_text = line.lstrip("- ").strip()
            temp_text = current_text + "\n- " + bullet_text if current_text else "- " + bullet_text
            
            # Check if adding this bullet exceeds the slide's capacity
            if not fits_on_slide(temp_text):
                # Save current content and start a new slide with the same heading
                text_frame = current_slide.shapes[1].text_frame
                text_frame.text = current_text.strip()
                current_slide = presentation.slides.add_slide(content_slide_layout)
                current_slide.shapes.title.text = sections["content"][sections["content"].index(line) - 1]  # Repeat heading
                current_text = "- " + bullet_text
            else:
                current_text = temp_text

    # Save the last slide's content
    if current_slide and current_text:
        text_frame = current_slide.shapes[1].text_frame
        text_frame.text = current_text.strip()

    # Add conclusion slide
    if "conclusion" in sections and sections["conclusion"]:
        slide = presentation.slides.add_slide(content_slide_layout)
        slide.shapes.title.text = "Conclusion"
        text_frame = slide.shapes[1].text_frame
        current_text = ""
        for line in sections["conclusion"]:
            if line.strip():
                temp_text = current_text + "\n" + line if current_text else line
                if not fits_on_slide(temp_text):
                    text_frame.text = current_text.strip()
                    slide = presentation.slides.add_slide(content_slide_layout)
                    slide.shapes.title.text = "Conclusion (cont.)"
                    text_frame = slide.shapes[1].text_frame
                    current_text = line
                else:
                    current_text = temp_text
        if current_text:
            text_frame.text = current_text.strip()

    # Save the presentation
    output_filename = "outputs/generated_presentation.pptx"
    presentation.save(output_filename)
    print(f"Presentation saved as {output_filename}!")

# Generate the presentation from the input text
# generate_presentation(input_text)